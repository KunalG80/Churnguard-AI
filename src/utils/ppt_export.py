"""
PowerPoint export — 5-slide boardroom deck.
Fixed:
  - KPI keys now match executive_summary() output
  - SHAP slide looks for shap_importance.png (saved by shap_explain with both names)
"""
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from src.config import PPT_PATH, FIGURES_DIR

logger = logging.getLogger(__name__)

DARK_BLUE  = RGBColor(0x1E, 0x3A, 0x5F)
LIGHT_GREY = RGBColor(0xF1, 0xF5, 0xF9)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
RED        = RGBColor(0xDC, 0x26, 0x26)
SLATE      = RGBColor(0x64, 0x74, 0x8B)


def _heading(slide, text: str, y: float = 0.2):
    txb = slide.shapes.add_textbox(Inches(0.3), Inches(y), Inches(9.4), Inches(0.7))
    tf  = txb.text_frame; tf.text = text
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(22); run.font.bold = True; run.font.color.rgb = DARK_BLUE


def _add_title_slide(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BLUE

    txb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf  = txb.text_frame
    tf.text = "ChurnGuard AI"
    tf.paragraphs[0].runs[0].font.size      = Pt(40)
    tf.paragraphs[0].runs[0].font.bold      = True
    tf.paragraphs[0].runs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    sub = tf.add_paragraph()
    sub.text = "Customer Retention Capital Allocation — Executive Briefing"
    sub.runs[0].font.size      = Pt(16)
    sub.runs[0].font.color.rgb = RGBColor(0xB0, 0xC4, 0xDE)
    sub.alignment = PP_ALIGN.CENTER


def _add_kpi_slide(prs, summary):
    """Fixed: KPI keys now match executive_summary() output."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, "Headline KPIs", y=0.3)

    kpis = [
        ("Revenue at Risk",
         f"₹{summary.get('Revenue at Risk', 0):,.0f}", RED),
        ("Total Recoverable Revenue",
         f"₹{summary.get('Total Recoverable Revenue', 0):,.0f}", GREEN),
        ("Net Value Created",
         f"₹{summary.get('Net Value Created', 0):,.0f}", GREEN),
        ("Overall ROI",
         f"{summary.get('Overall ROI', 0)*100:.1f}%", DARK_BLUE),
    ]
    for i, (label, value, colour) in enumerate(kpis):
        col = i % 2; row = i // 2
        box = slide.shapes.add_textbox(
            Inches(0.5 + col * 4.8), Inches(1.4 + row * 2.3),
            Inches(4.2), Inches(1.8),
        )
        tf = box.text_frame; tf.word_wrap = True
        p1 = tf.paragraphs[0]; p1.text = value
        p1.runs[0].font.size      = Pt(28)
        p1.runs[0].font.bold      = True
        p1.runs[0].font.color.rgb = colour
        p2 = tf.add_paragraph(); p2.text = label
        p2.runs[0].font.size      = Pt(12)
        p2.runs[0].font.color.rgb = SLATE


def _add_segment_slide(prs, seg_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, "Segment ROI Strategy")

    cols  = ["Risk_Tier", "Customers", "Revenue_at_Risk",
             "Retention_Investment", "Net_Value_Created", "Recommendation"]
    avail = [c for c in cols if c in seg_df.columns]
    rows_data = [[c.replace("_", " ") for c in avail]] + [
        [f"₹{v:,.0f}" if isinstance(v, float) else str(v)
         for v in seg_df[avail].iloc[i]]
        for i in range(len(seg_df))
    ]
    n_rows, n_cols = len(rows_data), len(avail)
    table = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.2), Inches(1.2),
        Inches(9.0), Inches(0.5) * n_rows,
    ).table

    for c_idx, header in enumerate(rows_data[0]):
        cell = table.cell(0, c_idx)
        cell.text = header
        cell.fill.solid(); cell.fill.fore_color.rgb = DARK_BLUE
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = WHITE; run.font.bold = True; run.font.size = Pt(9)

    for r_idx, row in enumerate(rows_data[1:], 1):
        bg = LIGHT_GREY if r_idx % 2 == 0 else WHITE
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(9)
            if avail[c_idx] == "Recommendation":
                run.font.color.rgb = GREEN if val == "INVEST" else RED
                run.font.bold = True


def _add_chart_slide(prs, fig_name: str, title: str):
    """Fixed: gracefully skips if chart not found, logs warning."""
    fig_path = FIGURES_DIR / fig_name
    if not fig_path.exists():
        logger.warning("Chart not found, skipping slide '%s': %s", title, fig_path)
        return
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _heading(slide, title)
    slide.shapes.add_picture(str(fig_path), Inches(1), Inches(1.3), width=Inches(8))


def export_ppt(bundle: dict):
    summary = bundle["summary"]
    seg_df  = bundle["segmentation"]

    PPT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, summary)
    _add_kpi_slide(prs, summary)
    _add_segment_slide(prs, seg_df)
    _add_chart_slide(prs, "revenue_by_segment.png", "Revenue Exposure by Segment")
    _add_chart_slide(prs, "shap_importance.png",    "Top Churn Drivers (SHAP)")

    try:
        prs.save(str(PPT_PATH))
        logger.info("PPT saved → %s", PPT_PATH)
    except Exception as e:
        logger.error("PPT save failed: %s", e); raise
